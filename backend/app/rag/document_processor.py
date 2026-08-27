"""
Docling-powered document processor for DeepTutor RAG pipeline.

Docling (IBM Research) is the primary parser — it provides:
  - Structure-aware PDF parsing (correct reading order, multi-column)
  - Native section heading hierarchy  
  - Table extraction as structured markdown
  - Formula/math preservation
  - OCR support for scanned PDFs
  - Image caption extraction
  - Clean markdown output

Falls back gracefully to pdfplumber → pypdf → PyPDF2 if Docling is
unavailable or conversion fails.

Architecture:
  process_document(file_path)
    └─► try_docling()       → DoclingDocument → DoclingChunker
    └─► try_pdfplumber()    → raw text → SemanticChunker  (fallback 1)
    └─► try_pypdf()         → raw text → SemanticChunker  (fallback 2)
    └─► try_pypdf2()        → raw text → SemanticChunker  (fallback 3)
"""
import os
import re
from pathlib import Path
from typing import List, Dict, Optional

from app.core.config import get_settings
from app.rag.topic_sanitizer import clean_and_format_topic, deduplicate_and_rank_topics, is_valid_academic_topic

settings = get_settings()

# Windows: suppress PyTorch MSVC compiler check error
os.environ.setdefault("DOCLING_INFERENCE_COMPILE_TORCH_MODELS", "false")

# ── Load chunker (for fallback path) ──────────────────────────────────────────
try:
    from app.rag.chunking_strategies import get_chunker
    _chunker = get_chunker(settings.CHUNKING_STRATEGY)
except Exception:
    _chunker = None


# ══════════════════════════════════════════════════════════════════════════════
# Docling converter (singleton, lazy-loaded)
# ══════════════════════════════════════════════════════════════════════════════
_docling_converter = None
_docling_available = None  # None = not yet checked


def _get_docling_converter():
    """
    Lazy-load and cache the Docling DocumentConverter.
    Returns converter or None if Docling is not installed.
    """
    global _docling_converter, _docling_available

    if _docling_available is False:
        return None
    if _docling_converter is not None:
        return _docling_converter

    try:
        import torch
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions, AcceleratorOptions
        from docling.document_converter import DocumentConverter, PdfFormatOption

        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_ocr = settings.DOCLING_ENABLE_OCR

        accel_opts = None
        # CUDA GPU acceleration for Docling layout & table vision models
        if torch.cuda.is_available():
            try:
                torch.cuda.set_device(0)
            except Exception:
                pass
            device_name = torch.cuda.get_device_name(0)
            print(f"[DOCLING GPU] CUDA enabled on {device_name}! Routing layout models to GPU...")
            accel_opts = AcceleratorOptions(num_threads=8, device="cuda:0")
            pipeline_options.accelerator_options = accel_opts
        else:
            print("[DOCLING CPU] CUDA not detected. Docling running on CPU.")

        format_options = {
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
        }

        # Enable Docling builtin Image OCR format options if supported
        try:
            from docling.datamodel.pipeline_options import ImagePipelineOptions
            from docling.document_converter import ImageFormatOption
            image_pipeline_options = ImagePipelineOptions()
            image_pipeline_options.do_ocr = True
            if accel_opts:
                image_pipeline_options.accelerator_options = accel_opts
            format_options[InputFormat.IMAGE] = ImageFormatOption(pipeline_options=image_pipeline_options)
        except Exception:
            pass

        _docling_converter = DocumentConverter(format_options=format_options)
        _docling_available = True
        return _docling_converter

    except Exception as e:
        print(f"[DOCLING WARN] Converter initialization error: {e}")
        _docling_available = False
        return None


# ══════════════════════════════════════════════════════════════════════════════
# Docling Chunker — exploits DoclingDocument native structure
# ══════════════════════════════════════════════════════════════════════════════
class DoclingChunker:
    """
    Chunks a DoclingDocument using its native hierarchy:
      - Iterates over document items (headings, paragraphs, tables, figures)
      - Groups paragraphs under their parent heading
      - Respects token budget per chunk
      - Produces rich metadata including section_title, page, item_type
    """

    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        self.chunk_size = chunk_size or settings.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or settings.CHUNK_OVERLAP

    def chunk_document(self, doc, source_name: str) -> List[Dict]:
        """
        doc: docling.datamodel.document.DoclingDocument
        Returns list of chunk dicts with text + metadata.
        """
        chunks: List[Dict] = []
        current_section = ""
        current_level = 0
        buffer: List[str] = []
        buffer_tokens = 0
        current_page = 1
        chunk_idx = 0

        def flush_buffer(clear_overlap=False) -> None:
            nonlocal buffer, buffer_tokens, chunk_idx
            text = " ".join(buffer).strip()
            if len(text) >= settings.MIN_CHUNK_CHARS:
                contextual_text = f"[Section: {current_section}]\n{text}" if current_section else text
                chunks.append({
                    "text": contextual_text,
                    "metadata": {
                        "source": source_name,
                        "page": current_page,
                        "section_title": current_section,
                        "section_level": current_level,
                        "chunk_index": chunk_idx,
                        "estimated_tokens": len(contextual_text) // 4,
                        "char_count": len(contextual_text),
                        "chunk_type": "docling",
                    }
                })
                chunk_idx += 1
            
            if clear_overlap:
                buffer = []
                buffer_tokens = 0
            else:
                # Keep overlap for next chunk
                overlap_text = " ".join(buffer)
                overlap_words = overlap_text.split()[-self.chunk_overlap:]
                buffer = overlap_words
                buffer_tokens = sum(len(w) for w in buffer) // 4

        try:
            for item, _ in doc.iterate_items():
                item_type = type(item).__name__

                if hasattr(item, 'prov') and item.prov:
                    current_page = item.prov[0].page_no

                if item_type in ("SectionHeaderItem", "TitleItem"):
                    flush_buffer(clear_overlap=True)
                    current_section = item.text.strip() if hasattr(item, 'text') else ""
                    current_level = getattr(item, 'level', 1)

                elif item_type in ("TextItem", "ParagraphItem"):
                    text = item.text.strip() if hasattr(item, 'text') else ""
                    if text:
                        text = re.sub(r'\s+', ' ', text)
                        token_est = len(text) // 4
                        if buffer_tokens + token_est > self.chunk_size and buffer:
                            flush_buffer()
                        buffer.append(text)
                        buffer_tokens += token_est

                elif item_type == "TableItem":
                    flush_buffer(clear_overlap=True)
                    table_md = item.export_to_markdown() if hasattr(item, 'export_to_markdown') else item.text
                    if table_md:
                        header = f"[Section: {current_section}]\n" if current_section else ""
                        chunks.append({
                            "text": f"{header}[TABLE]\n{table_md.strip()}",
                            "metadata": {"source": source_name, "page": current_page, "section_title": current_section, "chunk_index": chunk_idx, "chunk_type": "docling_table"}
                        })
                        chunk_idx += 1

                elif item_type == "FigureItem":
                    caption = item.caption.strip() if hasattr(item, 'caption') and item.caption else (item.text.strip() if hasattr(item, 'text') else "")
                    if caption:
                        buffer.append(f"[Figure: {caption}]")
                        buffer_tokens += len(caption) // 4

                elif item_type in ("EquationItem", "FormulaItem"):
                    formula = item.text.strip() if hasattr(item, 'text') else ""
                    if formula:
                        buffer.append(f"[Formula: {formula}]")
                        buffer_tokens += len(formula) // 4

                elif item_type == "ListItem":
                    text = item.text.strip() if hasattr(item, 'text') else ""
                    if text:
                        buffer.append(f"\u2022 {text}")
                        buffer_tokens += len(text) // 4

                elif item_type == "CodeItem":
                    code = item.text.strip() if hasattr(item, 'text') else ""
                    if code and len(code) >= 10:
                        flush_buffer(clear_overlap=True)
                        prefix = f"[Section: {current_section}]\n" if current_section else ""
                        code_text = f"{prefix}[CODE]\n```\n{code}\n```"
                        chunks.append({
                            "text": code_text,
                            "metadata": {
                                "source": source_name, "page": current_page,
                                "section_title": current_section, "section_level": current_level,
                                "chunk_index": chunk_idx, "estimated_tokens": len(code_text) // 4,
                                "char_count": len(code_text), "chunk_type": "docling_code",
                            }
                        })
                        chunk_idx += 1

        except Exception:
            pass

        # Flush remaining buffer
        if buffer:
            flush_buffer()

        return chunks


# ══════════════════════════════════════════════════════════════════════════════
# Docling extraction path
# ══════════════════════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════════════════════
# Docling extraction path
# ══════════════════════════════════════════════════════════════════════════════
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError

def _run_docling_conversion(file_path: str) -> List[Dict]:
    """Internal helper to convert document with Docling."""
    converter = _get_docling_converter()
    if converter is None:
        return []
    source_name = Path(file_path).name
    result = converter.convert(file_path)
    doc = result.document

    chunker = DoclingChunker(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
    )
    chunks = chunker.chunk_document(doc, source_name)

    if not chunks:
        md_text = doc.export_to_markdown()
        if md_text and _chunker is not None:
            return _chunker.chunk(md_text, {
                "source": source_name,
                "page": 1,
                "file_path": file_path,
            })
    return chunks


def _try_docling(file_path: str) -> List[Dict]:
    """
    Primary parser: uses Docling for structure-aware PDF extraction with strict timeout.
    Returns list of chunks or [] on failure/timeout.
    """
    if not getattr(settings, "ENABLE_DOCLING", False):
        return []

    timeout_sec = getattr(settings, "DOCLING_TIMEOUT_SECONDS", 12)
    try:
        with ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(_run_docling_conversion, file_path)
            return future.result(timeout=timeout_sec)
    except (FuturesTimeoutError, Exception) as e:
        # Fall back gracefully to fast text parsers (pdfplumber/pypdf)
        return []


# ══════════════════════════════════════════════════════════════════════════════
# Fallback parsers (pdfplumber → pypdf → pypdfium2 → PyPDF2)
# ══════════════════════════════════════════════════════════════════════════════
def _legacy_split(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Character-based fallback splitter."""
    char_size = chunk_size * 4
    char_overlap = overlap * 4
    chunks, start = [], 0
    while start < len(text):
        chunk = text[start:start + char_size].strip()
        if chunk and len(chunk) >= settings.MIN_CHUNK_CHARS:
            chunks.append(chunk)
        start += char_size - char_overlap
    return chunks


def _chunk_page_text(text: str, page_num: int, file_path: str) -> List[Dict]:
    """Chunk a single page of text using configured chunker."""
    if not text or len(text) < settings.MIN_CHUNK_CHARS:
        return []
    source_name = Path(file_path).name
    meta = {"source": source_name, "page": page_num, "file_path": file_path}

    if _chunker is not None:
        return _chunker.chunk(text, meta)

    # Absolute fallback
    return [
        {
            "text": c,
            "metadata": {**meta, "chunk_index": i, "section_title": "",
                         "section_level": 0, "chunk_type": "legacy",
                         "estimated_tokens": len(c)//4, "char_count": len(c)},
        }
        for i, c in enumerate(_legacy_split(text, settings.CHUNK_SIZE, settings.CHUNK_OVERLAP))
    ]


def _try_pdfplumber(file_path: str) -> List[Dict]:
    try:
        import pdfplumber
        chunks = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = re.sub(r'\s+', ' ', page.extract_text() or "").strip()
                chunks.extend(_chunk_page_text(text, page_num, file_path))
        return chunks
    except Exception:
        return []


def _try_pypdf(file_path: str) -> List[Dict]:
    try:
        import pypdf
        chunks = []
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = re.sub(r'\s+', ' ', page.extract_text() or "").strip()
                chunks.extend(_chunk_page_text(text, page_num, file_path))
        return chunks
    except Exception:
        return []


def _try_pypdfium2(file_path: str) -> List[Dict]:
    try:
        import pypdfium2 as pdfium
        chunks = []
        pdf = pdfium.PdfDocument(file_path)
        for page_num, page in enumerate(pdf, 1):
            textpage = page.get_textpage()
            text = re.sub(r'\s+', ' ', textpage.get_text_range() or "").strip()
            chunks.extend(_chunk_page_text(text, page_num, file_path))
        return chunks
    except Exception:
        return []


def _try_pypdf2(file_path: str) -> List[Dict]:
    try:
        import PyPDF2
        chunks = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = re.sub(r'\s+', ' ', page.extract_text() or "").strip()
                chunks.extend(_chunk_page_text(text, page_num, file_path))
        return chunks
    except Exception:
        return []


# ══════════════════════════════════════════════════════════════════════════════
# OCR Engine Helpers (EasyOCR + PyTesseract + Docling)
# ══════════════════════════════════════════════════════════════════════════════
_easyocr_reader = None

def _get_easyocr_reader():
    global _easyocr_reader
    if _easyocr_reader is not None:
        return _easyocr_reader
    try:
        import easyocr
        import torch
        use_gpu = torch.cuda.is_available()
        print(f"[OCR] Initializing EasyOCR engine (GPU={use_gpu})...")
        _easyocr_reader = easyocr.Reader(['en'], gpu=use_gpu, verbose=False)
        return _easyocr_reader
    except Exception as e:
        print(f"[OCR WARN] EasyOCR initialization note: {e}")
        return None



def _try_gemini_vlm(file_path_or_img, prompt: Optional[str] = None) -> str:
    """Run Google Gemini 2.0 / 1.5 Flash VLM to transcribe an image or document page."""
    if not getattr(settings, "ENABLE_VLM_PARSER", True):
        return ""
    try:
        from app.rag.gemini_client import gemini_client
        api_key = os.environ.get("GEMINI_API_KEY") or getattr(settings, "GEMINI_API_KEY", "")
        if not api_key or len(api_key.strip()) < 5:
            return ""

        target_vlm_model = getattr(settings, "GEMINI_VLM_MODEL", "gemini-2.0-flash") or "gemini-2.0-flash"
        text = gemini_client.sync_transcribe_image_vlm(
            file_path_or_img,
            prompt=prompt,
            model=target_vlm_model
        )
        if text and len(text.strip()) >= 10:
            return text.strip()
    except Exception as e:
        print(f"[VLM WARN] Gemini VLM transcription error: {e}")
    return ""


def _ocr_extract_text(file_path_or_img) -> str:
    """Run Gemini Flash VLM -> EasyOCR -> PyTesseract OCR fallback chain on an image file or PIL image."""
    # 1. Primary: Gemini Flash VLM
    vlm_text = _try_gemini_vlm(file_path_or_img)
    if vlm_text:
        return vlm_text

    # 2. EasyOCR
    reader = _get_easyocr_reader()
    if reader is not None:
        try:
            results = reader.readtext(file_path_or_img, detail=0)
            text = " ".join(results).strip()
            if text and len(text) >= 10:
                return text
        except Exception as e:
            print(f"[OCR WARN] EasyOCR extraction error: {e}")

    # 3. PyTesseract
    try:
        import pytesseract
        from PIL import Image
        if isinstance(file_path_or_img, (str, Path)):
            img = Image.open(file_path_or_img)
        else:
            img = file_path_or_img
        text = pytesseract.image_to_string(img).strip()
        if text and len(text) >= 10:
            return text
    except Exception as e:
        print(f"[OCR WARN] PyTesseract extraction error: {e}")

    return ""


def _ocr_scanned_pdf(file_path: str) -> List[Dict]:
    """Render scanned PDF pages as images and extract text using Gemini Flash VLM / OCR."""
    chunks = []
    source_name = Path(file_path).name
    max_pages = getattr(settings, "VLM_MAX_PAGES_PER_DOC", 50) or 50
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(file_path)
        total_pages = min(len(pdf), max_pages)

        for page_num in range(1, total_pages + 1):
            page = pdf[page_num - 1]
            pil_image = page.render(scale=2.0).to_pil()
            text = _ocr_extract_text(pil_image)
            if text and len(text) >= 15:
                page_chunks = _chunk_page_text(text, page_num, file_path)
                chunks.extend(page_chunks)
    except Exception as e:
        print(f"[OCR WARN] Scanned PDF VLM/OCR error for {file_path}: {e}")
    return chunks



# ══════════════════════════════════════════════════════════════════════════════
# Public API
# ══════════════════════════════════════════════════════════════════════════════
def process_pdf(file_path: str) -> List[Dict]:
    """
    Extract and chunk a PDF using high-speed cascading parser priority:
      1. pypdfium2 (sub-second fast text extraction)
      2. pdfplumber (layout & table extraction)
      3. pypdf (python text fallback)
      4. Scanned PDF OCR (renders pages & extracts text if PDF has no embedded text)
      5. Docling (structure-aware ML parser)
      6. PyPDF2 (legacy fallback)
    """
    # 1. pypdfium2
    chunks = _try_pypdfium2(file_path)
    if chunks and any(len(c.get("text", "").strip()) >= 50 for c in chunks):
        return chunks

    # 2. pdfplumber
    chunks = _try_pdfplumber(file_path)
    if chunks and any(len(c.get("text", "").strip()) >= 50 for c in chunks):
        return chunks

    # 3. pypdf
    chunks = _try_pypdf(file_path)
    if chunks and any(len(c.get("text", "").strip()) >= 50 for c in chunks):
        return chunks

    # 4. Scanned / Image-based PDF OCR
    print(f"[OCR] PDF {file_path} appears to be scanned or image-based. Running OCR page parser...")
    ocr_chunks = _ocr_scanned_pdf(file_path)
    if ocr_chunks:
        return ocr_chunks

    # 5. Docling
    if getattr(settings, "ENABLE_DOCLING", False):
        chunks = _try_docling(file_path)
        if chunks:
            return chunks

    # 6. PyPDF2
    return _try_pypdf2(file_path)



def process_docx(file_path: str) -> List[Dict]:
    """Extract text, headers, and tables from Word (.docx) documents."""
    try:
        import docx
        doc = docx.Document(file_path)
        source_name = Path(file_path).name
        full_text_blocks = []
        current_section = ""

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            if p.style and p.style.name.startswith("Heading"):
                current_section = text
                full_text_blocks.append(f"\n# {text}\n")
            else:
                full_text_blocks.append(text)

        # Extract tables as markdown tables
        for table in doc.tables:
            table_lines = []
            for i, row in enumerate(table.rows):
                row_vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                table_lines.append("| " + " | ".join(row_vals) + " |")
                if i == 0:
                    table_lines.append("| " + " | ".join(["---"] * len(row_vals)) + " |")
            if table_lines:
                full_text_blocks.append("\n" + "\n".join(table_lines) + "\n")

        joined_text = "\n".join(full_text_blocks).strip()
        if not joined_text:
            return []

        meta = {"source": source_name, "page": 1, "section_title": current_section, "file_path": file_path}
        chunks = _chunker.chunk(joined_text, meta) if _chunker is not None else []
        if not chunks:
            chunks = _chunk_page_text(joined_text, 1, file_path)
        return chunks
    except Exception as e:
        # Fallback to plain text read if docx library fails
        return process_txt(file_path)


def process_csv(file_path: str) -> List[Dict]:
    """Extract and format tabular CSV data as structured Markdown tables."""
    try:
        import csv
        source_name = Path(file_path).name
        chunks = []
        rows = []

        with open(file_path, mode="r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                if row:
                    rows.append(row)

        if not rows:
            return []

        header = rows[0]
        data_rows = rows[1:]

        # Batch 40 rows per chunk to preserve tabular context
        batch_size = 40
        for i in range(0, max(1, len(data_rows)), batch_size):
            batch = data_rows[i:i + batch_size]
            table_lines = [
                f"# CSV Dataset: {source_name} (Rows {i+1} to {i+len(batch)})",
                "| " + " | ".join([str(h).strip() for h in header]) + " |",
                "| " + " | ".join(["---"] * len(header)) + " |"
            ]
            for row in batch:
                row_str = "| " + " | ".join([str(val).strip().replace("\n", " ") for val in row]) + " |"
                table_lines.append(row_str)

            chunk_text = "\n".join(table_lines)
            page_num = (i // batch_size) + 1
            meta = {
                "source": source_name,
                "page": page_num,
                "section_title": f"Rows {i+1}-{i+len(batch)}",
                "file_path": file_path
            }

            chunks.append({
                "text": chunk_text,
                "metadata": {
                    **meta,
                    "chunk_index": len(chunks),
                    "chunk_type": "csv_table",
                    "estimated_tokens": len(chunk_text) // 4,
                    "char_count": len(chunk_text)
                }
            })

        return chunks
    except Exception as e:
        return process_txt(file_path)


def process_excel(file_path: str) -> List[Dict]:
    """Extract Excel spreadsheets (.xlsx, .xls) as Markdown table chunks per worksheet."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        source_name = Path(file_path).name
        chunks = []

        for page_num, sheet_name in enumerate(wb.sheetnames, 1):
            sheet = wb[sheet_name]
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            clean_rows = []
            for r in rows:
                if any(cell is not None for cell in r):
                    clean_rows.append([str(c).strip() if c is not None else "" for c in r])

            if not clean_rows:
                continue

            header = clean_rows[0]
            table_lines = [
                f"# Excel Sheet: {sheet_name}",
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(["---"] * len(header)) + " |"
            ]
            for r in clean_rows[1:]:
                table_lines.append("| " + " | ".join(r) + " |")

            chunk_text = "\n".join(table_lines)
            meta = {
                "source": source_name,
                "page": page_num,
                "section_title": f"Sheet: {sheet_name}",
                "file_path": file_path
            }
            chunks.append({
                "text": chunk_text,
                "metadata": {
                    **meta,
                    "chunk_index": len(chunks),
                    "chunk_type": "excel_sheet",
                    "estimated_tokens": len(chunk_text) // 4,
                    "char_count": len(chunk_text)
                }
            })
        return chunks
    except Exception:
        return process_txt(file_path)


def process_pptx(file_path: str) -> List[Dict]:
    """Extract PowerPoint presentation (.pptx) slides into formatted slide chunks."""
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        source_name = Path(file_path).name
        chunks = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        if shape == slide.shapes.title:
                            title = txt
                        else:
                            slide_texts.append(txt)

            slide_body = "\n\n".join(slide_texts)
            chunk_text = f"# Slide {slide_num}: {title}\n\n{slide_body}".strip()
            if len(chunk_text) >= 10:
                meta = {
                    "source": source_name,
                    "page": slide_num,
                    "section_title": f"Slide {slide_num}: {title}",
                    "file_path": file_path
                }
                chunks.append({
                    "text": chunk_text,
                    "metadata": {
                        **meta,
                        "chunk_index": len(chunks),
                        "chunk_type": "pptx_slide",
                        "estimated_tokens": len(chunk_text) // 4,
                        "char_count": len(chunk_text)
                    }
                })
        return chunks
    except Exception:
        return process_txt(file_path)


def process_html(file_path: str) -> List[Dict]:
    """Extract clean text and structure from HTML pages."""
    try:
        from bs4 import BeautifulSoup
        raw_html = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(raw_html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n").strip()
        text = re.sub(r'\n{3,}', '\n\n', text)
        meta = {"source": Path(file_path).name, "page": 1, "file_path": file_path}
        if _chunker is not None:
            return _chunker.chunk(text, meta)
        return _chunk_page_text(text, 1, file_path)
    except Exception:
        return process_txt(file_path)


def process_json(file_path: str) -> List[Dict]:
    """Parse JSON data files into formatted structured text chunks."""
    try:
        import json
        raw = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        data = json.loads(raw)
        formatted_text = json.dumps(data, indent=2)
        meta = {"source": Path(file_path).name, "page": 1, "file_path": file_path}
        if _chunker is not None:
            return _chunker.chunk(formatted_text, meta)
        return _chunk_page_text(formatted_text, 1, file_path)
    except Exception:
        return process_txt(file_path)


def process_txt(file_path: str) -> List[Dict]:
    """Read text/markdown file and split into chunks."""
    try:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r'\s+', ' ', text).strip()
        if not text:
            return []
        meta = {"source": Path(file_path).name, "page": 1, "file_path": file_path}
        if _chunker is not None:
            return _chunker.chunk(text, meta)
        return [
            {
                "text": c,
                "metadata": {**meta, "chunk_index": i, "section_title": "",
                             "section_level": 0, "chunk_type": "legacy",
                             "estimated_tokens": len(c)//4, "char_count": len(c)},
            }
            for i, c in enumerate(_legacy_split(text, settings.CHUNK_SIZE, settings.CHUNK_OVERLAP))
        ]
    except Exception as e:
        raise RuntimeError(f"Text processing failed: {e}")


def process_image(file_path: str) -> List[Dict]:
    """
    Extract text, diagram labels, equations, and structure from images (.png, .jpg, .jpeg, .webp, .bmp, .tiff)
    using EasyOCR / PyTesseract / Docling OCR engines.
    """
    source_name = Path(file_path).name

    # 1. Run EasyOCR / PyTesseract OCR extraction
    text = _ocr_extract_text(file_path)

    # 2. Docling OCR fallback
    if not text or len(text) < 10:
        docling_chunks = _try_docling(file_path)
        if docling_chunks:
            return docling_chunks

    if text and len(text) >= 10:
        meta = {
            "source": source_name,
            "page": 1,
            "section_title": f"Image Diagram: {source_name}",
            "file_path": file_path,
        }
        formatted_text = text if text.startswith("#") else f"# Image Document / Diagram: {source_name}\n\n{text}"
        return _chunk_page_text(formatted_text, 1, file_path) or [{
            "text": formatted_text,
            "metadata": {
                **meta,
                "chunk_index": 0,
                "chunk_type": "vlm_image" if getattr(settings, "ENABLE_VLM_PARSER", True) else "image_ocr",
                "estimated_tokens": len(formatted_text) // 4,
                "char_count": len(formatted_text),
            }
        }]


    meta = {"source": source_name, "page": 1, "file_path": file_path}
    img_text = f"[Image Document: {source_name}]\nImage uploaded. No clear text labels found by OCR."
    return [{
        "text": img_text,
        "metadata": {
            **meta,
            "chunk_index": 0,
            "chunk_type": "image_ocr",
            "estimated_tokens": len(img_text) // 4,
            "char_count": len(img_text),
        }
    }]



def process_document(file_path: str) -> List[Dict]:
    """Auto-detect file type and process with format-specific parsers."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return process_pdf(file_path)
    elif ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}:
        return process_image(file_path)
    elif ext in {".docx", ".doc"}:
        return process_docx(file_path)
    elif ext == ".csv":
        return process_csv(file_path)
    elif ext in {".xlsx", ".xls"}:
        return process_excel(file_path)
    elif ext in {".pptx", ".ppt"}:
        return process_pptx(file_path)
    elif ext in {".html", ".htm"}:
        return process_html(file_path)
    elif ext == ".json":
        return process_json(file_path)
    elif ext in {".txt", ".md", ".rst", ".log", ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp"}:
        return process_txt(file_path)
    else:
        # Fallback to plain text parser
        return process_txt(file_path)


def is_docling_available() -> bool:
    """Returns True if Docling is installed and usable."""
    return _get_docling_converter() is not None


def get_parser_info() -> Dict:
    """Return info about which parsers are available."""
    docling_ok = is_docling_available()
    pdfplumber_ok = True
    try:
        import pdfplumber
    except ImportError:
        pdfplumber_ok = False

    return {
        "primary": "docling" if docling_ok else "pdfplumber",
        "docling": docling_ok,
        "docling_version": _get_docling_version() if docling_ok else None,
        "ocr_enabled": settings.DOCLING_ENABLE_OCR,
        "pdfplumber": pdfplumber_ok,
        "chunking_strategy": settings.CHUNKING_STRATEGY,
    }


def _get_docling_version() -> Optional[str]:
    try:
        import docling
        return getattr(docling, "__version__", "unknown")
    except Exception:
        return None


def extract_key_topics(chunks: List[Dict]) -> List[str]:
    """
    Extract high-value academic concepts, headings, algorithms, and technical topics
    from document chunks. Filters out publisher metadata, country names, database tags, and noise.
    """
    if not chunks:
        return []

    META_STOPWORDS = {
        "the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with",
        "by", "from", "up", "about", "into", "over", "after", "is", "are", "was", "were",
        "be", "been", "being", "have", "has", "had", "do", "does", "did", "this", "that",
        "these", "those", "it", "its", "page", "pages", "pdf", "figure", "table", "chapter",
        "section", "author", "authors", "editor", "volume", "issue", "journal", "abstract",
        "introduction", "conclusion", "references", "http", "https", "doi", "isbn", "university",
        "department", "press", "rights", "reserved", "copyright", "edition", "published",
        "ieee", "south africa", "science core collection", "sci-expanded", "keywords plus",
        "web of science", "elsevier", "springer", "wiley", "taylor", "francis", "thomson reuters",
        "google scholar", "scopus", "proceedings", "conference", "symposium", "institution",
        "tc", "tp", "cpp", "lr", "roc", "usa", "uk", "china", "india", "japan", "germany"
    }

    ACRONYM_MAP = {
        "SVM": "Support Vector Machines (SVM)",
        "KNN": "K-Nearest Neighbors (KNN)",
        "RF": "Random Forest (RF)",
        "ML": "Machine Learning (ML)",
        "AI": "Artificial Intelligence (AI)",
        "CNN": "Convolutional Neural Networks (CNN)",
        "RNN": "Recurrent Neural Networks (RNN)",
        "LSTM": "Long Short-Term Memory (LSTM)",
        "BERT": "BERT Language Model",
        "LLM": "Large Language Models (LLM)",
        "NLP": "Natural Language Processing (NLP)",
        "PCA": "Principal Component Analysis (PCA)",
        "RAG": "Retrieval-Augmented Generation (RAG)",
    }

    candidates_counts: Dict[str, int] = {}

    for chunk in chunks:
        # 1. Section titles from Docling/pdfplumber metadata
        section_title = chunk.get("metadata", {}).get("section_title", "")
        if section_title:
            clean_sec = clean_and_format_topic(section_title)
            if clean_sec:
                candidates_counts[clean_sec] = candidates_counts.get(clean_sec, 0) + 6

        text = chunk.get("text", "")
        if not text:
            continue

        # 2. Text headings / numbered section headers
        for h in re.findall(r'(?:^|\n)(?:#{1,4}\s*|\d+(?:\.\d+)*\s+)?([A-Z][A-Za-z0-9\s\-\:\(\)]{3,50})(?=\n|\:|\.|\ {2,})', text):
            clean_h = clean_and_format_topic(h)
            if clean_h:
                candidates_counts[clean_h] = candidates_counts.get(clean_h, 0) + 4

        # 3. Capitalized multi-word concepts (e.g. 'Treaty of Vienna', 'Unification of Italy', 'Decision Trees')
        for c in re.findall(r'\b([A-Z][A-Za-z0-9\-]+(?:\s+(?:of|and|the|in|for|to)\s+[A-Z][A-Za-z0-9\-]+|\s+[A-Z][A-Za-z0-9\-]+){1,3})\b', text):
            clean_c = clean_and_format_topic(c)
            if clean_c:
                candidates_counts[clean_c] = candidates_counts.get(clean_c, 0) + 2

        # 4. Known domain acronyms
        for word in re.findall(r'\b([A-Z]{2,6})\b', text):
            if word in ACRONYM_MAP:
                mapped = ACRONYM_MAP[word]
                candidates_counts[mapped] = candidates_counts.get(mapped, 0) + 3

    sorted_candidates = [topic for topic, _ in sorted(candidates_counts.items(), key=lambda x: x[1], reverse=True)]
    return deduplicate_and_rank_topics(sorted_candidates, max_topics=15)
